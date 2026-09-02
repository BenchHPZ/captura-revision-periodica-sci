#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Prepara la plantilla que consume el informe fotográfico, a partir de la
plantilla corporativa 'Reporte sistemas.pptx'.

Hace dos cosas, ninguna de las dos disponible en la plantilla corporativa
tal cual viene:

1. Agrega una novena diapositiva con el layout 'Elemento', que sirve de
   molde para clonar una vez por cada elemento del ciclo. 'pptx-automizer',
   que es quien arma el informe, sólo sabe clonar diapositivas EXISTENTES
   (`addSlide(plantilla, numeroDeDiapositiva)`) — no tiene forma de crear
   una diapositiva a partir de un layout. Al molde se le RETIRAN todos los
   placeholders: al crear una diapositiva desde un layout se copian sus
   placeholders con el texto de ejemplo dentro ("Título en Deep Space
   Blue", "Click para editar el texto"…), que no es un texto guía sino
   contenido real y se vería impreso debajo de lo que dibuja el generador.
   Como el generador dibuja todo por su cuenta (ver docs/decisiones.md
   D-17), los placeholders no hacen falta: el fondo, el logo y los adornos
   vienen del layout, no de ellos.

2. Amplía la rejilla de la diapositiva de Agenda (5 casillas número+nombre,
   pensada para cinco sistemas fijos) a 10, para que el generador pueda
   listar entre 2 y 10 sistemas sin tocar la plantilla cada mes. Las 5
   casillas originales se conservan (se reposicionan y renombran); las 5
   que faltan se crean duplicando su XML tal cual —no redibujándolas—, así
   que el círculo, el color y la tipografía quedan idénticos a los de la
   plantilla real. Ver docs/decisiones.md D-17.

No modifica la plantilla corporativa: lee una y escribe otra.

Uso:
    python scripts/preparar_plantilla_informe.py \
        --origen "H:\\My Drive\\...\\Reporte sistemas.pptx" \
        --destino "H:\\My Drive\\...\\Plantilla_Informe.pptx"

Después hay que subirla al depósito:
    cd web && npx tsx scripts/subir-plantilla-informe.ts --archivo "<destino>" --confirmar
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8")

# Índice del layout 'Elemento' dentro de la plantilla corporativa. Los seis
# layouts son, en orden: Intro / Title with gradient 1 / Title with
# gradient 2 / 5x Agenda / Chapter divider / Elemento.
LAYOUT_ELEMENTO = 5
NOMBRE_LAYOUT_ESPERADO = "Elemento"
DIAPOSITIVAS_BASE_ESPERADAS = 8

NOMBRE_LAYOUT_AGENDA_ESPERADO = "5x Agenda Deep Space Blue"
NOMBRE_FORMA_NUMERO = "Text Placeholder 26"
NOMBRE_FORMA_NOMBRE = "Text Placeholder 27"
AGENDA_FILAS = 5
AGENDA_MAX_SISTEMAS = AGENDA_FILAS * 2


def preparar_agenda(slide) -> bool:
    """Amplía la rejilla número+nombre de la Agenda de 5 a
    `AGENDA_MAX_SISTEMAS` casillas, duplicando el XML de las 5 originales
    en vez de redibujarlas, para que el círculo, el color y la tipografía
    salgan idénticos a los de la plantilla real. Cada casilla final queda
    con un nombre único ('Agenda Numero N' / 'Agenda Nombre N') — las 5
    originales repetían el mismo nombre en las 5 y no se podían
    direccionar una por una. Devuelve False (y no toca nada) si la
    plantilla no trae la forma esperada."""
    from copy import deepcopy

    from pptx.oxml.ns import qn
    from pptx.util import Inches

    numeros = [sh for sh in slide.shapes if sh.name == NOMBRE_FORMA_NUMERO]
    nombres = [sh for sh in slide.shapes if sh.name == NOMBRE_FORMA_NOMBRE]
    if len(numeros) != AGENDA_FILAS or len(nombres) < AGENDA_FILAS:
        print(
            f"La diapositiva de Agenda trae {len(numeros)} casillas de número y "
            f"{len(nombres)} de nombre; se esperaban {AGENDA_FILAS} y al menos {AGENDA_FILAS}.",
            file=sys.stderr,
        )
        return False

    # Emparejar por la fila más cercana (el número y su nombre no
    # comparten exactamente el mismo 'top' en la plantilla real — están
    # a 0.05" de diferencia), antes de duplicar nada: una vez que haya
    # copias a la misma altura que el original, la posición deja de
    # alcanzar para distinguirlas. La plantilla trae, además, una casilla
    # de nombre vacía y duplicada exactamente sobre una real (fila 2,
    # columna 1) — de ahí que pueda haber más nombres que números; en un
    # empate de distancia se prefiere la que sí tiene texto, y la sobrante
    # se descarta al final.
    disponibles = list(nombres)
    originales = []
    for num in numeros:
        disponibles.sort(key=lambda n: (abs(n.top - num.top), 0 if n.text_frame.text.strip() else 1))
        pareja = disponibles.pop(0)
        originales.append((num, pareja))

    # Lo que sobre en 'disponibles' es un duplicado sin usar: se retira en
    # vez de dejarlo suelto en la diapositiva.
    for sobrante in disponibles:
        sobrante._element.getparent().remove(sobrante._element)

    def asignar_id(elemento, nuevo_id: int) -> None:
        cNvPr = elemento.find(f".//{qn('p:cNvPr')}")
        cNvPr.set("id", str(nuevo_id))

    siguiente_id = max(sh.shape_id for sh in slide.shapes) + 1
    duplicados = []
    for num_src, nom_src in originales:
        nuevo_num_el = deepcopy(num_src._element)
        nuevo_nom_el = deepcopy(nom_src._element)
        asignar_id(nuevo_num_el, siguiente_id)
        siguiente_id += 1
        asignar_id(nuevo_nom_el, siguiente_id)
        siguiente_id += 1
        num_src._element.addnext(nuevo_num_el)
        nom_src._element.addnext(nuevo_nom_el)
        duplicados.append((slide.shapes._shape_factory(nuevo_num_el), slide.shapes._shape_factory(nuevo_nom_el)))

    todas_las_parejas = originales + duplicados  # 10 parejas: 1-5 originales, 6-10 nuevas

    # El nombre no queda exactamente a la misma altura que su número en el
    # diseño original (~0.05" más arriba) — se conserva ese desfase en vez
    # de alinearlos, para no alterar el acomodo visual de la plantilla.
    desfase_nombre = originales[0][1].top - originales[0][0].top

    primera_fila_y = Inches(2.068)
    paso_fila_y = Inches(1.052)
    columnas_x = [(Inches(2.227), Inches(2.958)), (Inches(7.428), Inches(8.158))]

    for i, (num, nom) in enumerate(todas_las_parejas, start=1):
        fila = (i - 1) % AGENDA_FILAS
        columna = 0 if i <= AGENDA_FILAS else 1
        y = primera_fila_y + fila * paso_fila_y
        num.left, num.top = columnas_x[columna][0], y
        nom.left, nom.top = columnas_x[columna][1], y + desfase_nombre
        num.name = f"Agenda Numero {i}"
        nom.name = f"Agenda Nombre {i}"
        # El texto de ejemplo se deja tal cual, sin vaciarlo: vaciar un
        # <a:p> le quita su <a:r> (la corrida de texto), y sin una corrida
        # existente 'modifyElement' de pptx-automizer no tiene qué
        # reemplazar (mismo problema de los placeholders vacíos — ver
        # docs/decisiones.md D-17). No importa que quede un nombre de
        # sistema de ejemplo aquí: el generador llena las 10 casillas en
        # cada corrida, usadas o no, así que este texto nunca sobrevive a
        # un informe real.

    return True


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--origen", required=True, help="Ruta a 'Reporte sistemas.pptx'")
    ap.add_argument("--destino", required=True, help="Ruta del .pptx a escribir")
    args = ap.parse_args()

    origen = Path(args.origen)
    destino = Path(args.destino)

    if not origen.exists():
        print(f"No existe: {origen}", file=sys.stderr)
        return 1

    from pptx import Presentation

    prs = Presentation(str(origen))

    # Comprobaciones antes de escribir: si la plantilla corporativa cambia
    # de forma, es mejor detenerse aquí que producir un informe torcido.
    if len(prs.slide_layouts) <= LAYOUT_ELEMENTO:
        print(f"La plantilla sólo tiene {len(prs.slide_layouts)} layouts; se esperaba al menos {LAYOUT_ELEMENTO + 1}.", file=sys.stderr)
        return 1

    layout = prs.slide_layouts[LAYOUT_ELEMENTO]
    if layout.name != NOMBRE_LAYOUT_ESPERADO:
        print(f"El layout {LAYOUT_ELEMENTO} se llama {layout.name!r}, no {NOMBRE_LAYOUT_ESPERADO!r}.", file=sys.stderr)
        print("La plantilla corporativa cambió: revisar antes de seguir.", file=sys.stderr)
        return 1

    base = len(prs.slides)
    if base != DIAPOSITIVAS_BASE_ESPERADAS:
        print(f"Aviso: la plantilla trae {base} diapositivas base, no {DIAPOSITIVAS_BASE_ESPERADAS}.")
        print("El generador asume 1 intro + 1 portada + 1 agenda + un divisor genérico por sistema.")

    agenda = next((s for s in prs.slides if s.slide_layout.name == NOMBRE_LAYOUT_AGENDA_ESPERADO), None)
    if agenda is None:
        print(f"No se encontró ninguna diapositiva con el layout {NOMBRE_LAYOUT_AGENDA_ESPERADO!r}.", file=sys.stderr)
        return 1
    if not preparar_agenda(agenda):
        return 1

    molde = prs.slides.add_slide(layout)

    # Retirar los placeholders copiados del layout: traen el texto de
    # ejemplo dentro y saldría impreso debajo de lo que dibuja el
    # generador.
    retirados = []
    for ph in list(molde.placeholders):
        retirados.append(f"idx={ph.placeholder_format.idx} {ph.name!r}")
        ph._element.getparent().remove(ph._element)

    if molde.shapes:
        print(f"Aviso: quedaron {len(molde.shapes)} formas en el molde que no eran placeholders.")

    destino.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destino))

    print(f"Origen:  {origen}")
    print(f"Destino: {destino}")
    print(f"Diapositivas: {base} base + 1 molde de 'Elemento' = {base + 1}")
    print(f"Agenda ampliada a {AGENDA_MAX_SISTEMAS} casillas (número + nombre).")
    print(f"Placeholders retirados del molde: {len(retirados)}")
    for r in retirados:
        print(f"   - {r}")
    print("\nListo. Ahora súbela al depósito con:")
    print(f'  cd web && npx tsx scripts/subir-plantilla-informe.ts --archivo "{destino}" --confirmar')
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
